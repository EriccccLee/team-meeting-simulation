"""
POST /api/run       — 시뮬레이션 시작, session_id 즉시 반환
GET  /api/stream/{session_id} — SSE 이벤트 스트림

파일 전처리(PDF/Excel/Word/PowerPoint 변환)는 백그라운드 스레드에서 수행되며
preprocessing 이벤트로 진행 상황을 실시간 전달합니다.
"""
from __future__ import annotations

import asyncio
import io
import json
import queue as stdlib_queue
import sys
import tempfile
import uuid
from pathlib import Path
from typing import Annotated

from fastapi import APIRouter, File, Form, HTTPException, UploadFile
from fastapi.responses import StreamingResponse

# simulation 패키지가 상위 디렉터리에 있으므로 sys.path 조정
_ROOT = Path(__file__).parent.parent.parent
if str(_ROOT) not in sys.path:
    sys.path.insert(0, str(_ROOT))

from simulation.agents import MeetingAgent, ModeratorAgent
from simulation.model_client import ClaudeCodeModelClient
from simulation.orchestrator import MeetingOrchestrator, OrchestratorConfig
from simulation.session import MeetingSession
from simulation.cli import _load_agent_config, _load_file_contents

router = APIRouter()

# session_id → SimpleQueue (thread-safe, no asyncio needed)
_sessions: dict[str, stdlib_queue.SimpleQueue] = {}


# ── 파일 포맷 변환 헬퍼 ───────────────────────────────────────────────────────

def _excel_to_md(filename: str, raw: bytes) -> str:
    """xlsx/xls → 마크다운 표"""
    try:
        import pandas as pd
        dfs = pd.read_excel(io.BytesIO(raw), sheet_name=None)
        parts = []
        for sheet_name, df in dfs.items():
            parts.append(f"## {sheet_name}\n\n{df.to_markdown(index=False)}\n")
        return "\n".join(parts)
    except Exception as e:
        return f"[Excel 변환 실패: {e}]"


def _csv_to_md(filename: str, raw: bytes) -> str:
    """csv → 마크다운 표"""
    try:
        import pandas as pd
        try:
            df = pd.read_csv(io.BytesIO(raw), encoding="utf-8")
        except UnicodeDecodeError:
            df = pd.read_csv(io.BytesIO(raw), encoding="cp949")
        return df.to_markdown(index=False)
    except Exception:
        try:
            return raw.decode("utf-8")
        except UnicodeDecodeError:
            return raw.decode("cp949", errors="replace")


def _docx_to_md(filename: str, raw: bytes) -> str:
    """docx → 마크다운 텍스트"""
    try:
        from docx import Document
        doc = Document(io.BytesIO(raw))
        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name if para.style else ""
            if "Heading 1" in style:
                parts.append(f"# {text}")
            elif "Heading 2" in style:
                parts.append(f"## {text}")
            elif "Heading 3" in style:
                parts.append(f"### {text}")
            else:
                parts.append(text)
        return "\n\n".join(parts)
    except ImportError:
        return "[docx 지원: pip install python-docx]"
    except Exception as e:
        return f"[Word 문서 변환 실패: {e}]"


def _pptx_to_md(filename: str, raw: bytes) -> str:
    """pptx → 마크다운 텍스트 (슬라이드별)"""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw))
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            if slide_texts:
                parts.append(f"## 슬라이드 {i}\n\n" + "\n\n".join(slide_texts))
        return "\n\n---\n\n".join(parts) if parts else "[빈 프레젠테이션]"
    except ImportError:
        return "[pptx 지원: pip install python-pptx]"
    except Exception as e:
        return f"[PowerPoint 변환 실패: {e}]"


def _decode_text(raw: bytes) -> str:
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        return raw.decode("cp949", errors="replace")


# ── 시뮬레이션 실행 (동기, 별도 스레드) ──────────────────────────────────────

def _run_simulation(
    session_id: str,
    topic: str,
    participant_slugs: list[str],
    rounds: int,
    raw_files: list[tuple[str, bytes]],
) -> None:
    """파일 전처리 → orchestrator.run() 을 동기 스레드에서 실행.

    stdlib_queue.SimpleQueue 는 thread-safe 하므로 asyncio 없이 직접 put().
    SSE generator 쪽에서 get_nowait() 폴링으로 소비함.
    """
    q = _sessions[session_id]

    def emit(event: dict) -> None:
        q.put(event)  # SimpleQueue.put() is thread-safe

    try:
        # ── 1. 파일 전처리 (preprocessing SSE 이벤트 emit) ─────────────────
        file_contents: dict[str, str] = {}
        total = len(raw_files)

        for idx, (filename, raw) in enumerate(raw_files, start=1):
            ext = Path(filename).suffix.lower()
            # default-arg trick prevents late-binding closure issues
            def pre(msg: str, done: bool = False,
                    _fn=filename, _idx=idx, _tot=total) -> None:
                emit({"type": "preprocessing", "message": msg,
                      "filename": _fn, "index": _idx, "total": _tot, "done": done})

            if ext == ".pdf":
                pre(f"PDF 변환 중: {filename}")
                tmp_path = None
                try:
                    with tempfile.NamedTemporaryFile(
                        suffix=".pdf", delete=False, dir=str(_ROOT)
                    ) as tmp:
                        tmp.write(raw)
                        tmp_path = tmp.name
                    loaded = _load_file_contents([tmp_path])
                    file_contents.update(loaded)
                finally:
                    if tmp_path:
                        Path(tmp_path).unlink(missing_ok=True)
                pre(f"PDF 변환 완료: {filename}", done=True)

            elif ext in (".xlsx", ".xls"):
                pre(f"Excel 변환 중: {filename}")
                file_contents[filename] = _excel_to_md(filename, raw)
                pre(f"Excel 변환 완료: {filename}", done=True)

            elif ext == ".csv":
                pre(f"CSV 변환 중: {filename}")
                file_contents[filename] = _csv_to_md(filename, raw)
                pre(f"CSV 변환 완료: {filename}", done=True)

            elif ext == ".docx":
                pre(f"Word 문서 변환 중: {filename}")
                file_contents[filename] = _docx_to_md(filename, raw)
                pre(f"Word 문서 변환 완료: {filename}", done=True)

            elif ext == ".pptx":
                pre(f"PowerPoint 변환 중: {filename}")
                file_contents[filename] = _pptx_to_md(filename, raw)
                pre(f"PowerPoint 변환 완료: {filename}", done=True)

            else:
                # 텍스트 계열 (.md, .txt 등) — 전처리 이벤트 없음
                file_contents[filename] = _decode_text(raw)

        # ── 2. 시뮬레이션 실행 ────────────────────────────────────────────────
        model_client = ClaudeCodeModelClient()
        agents: list[MeetingAgent] = []
        for slug in participant_slugs:
            config = _load_agent_config(slug)
            agents.append(MeetingAgent(config, model_client))

        moderator = ModeratorAgent(model_client, participant_slugs)
        session = MeetingSession(
            topic,
            participant_slugs,
            output_dir=str(_ROOT / "outputs"),
            emit=emit,
        )
        cfg = OrchestratorConfig(phase2_rounds=rounds)
        orchestrator = MeetingOrchestrator(agents, moderator, session, cfg)
        result = orchestrator.run(topic, file_contents)

        emit({"type": "done", "output_file": str(result.output_file)})

    except Exception as e:
        emit({"type": "error", "message": str(e)})
    finally:
        q.put(None)  # sentinel — SSE generator 종료 신호


# ── API 엔드포인트 ─────────────────────────────────────────────────────────────

@router.post("/run")
async def run_simulation(
    topic: Annotated[str, Form()],
    participants: Annotated[list[str], Form()],
    rounds: Annotated[int, Form()] = 3,
    files: list[UploadFile] = File(default=[]),
) -> dict:
    """session_id 를 즉시 반환합니다. 파일 전처리 + 시뮬레이션은 백그라운드 스레드에서 실행됩니다."""
    session_id = str(uuid.uuid4())
    q: stdlib_queue.SimpleQueue = stdlib_queue.SimpleQueue()
    _sessions[session_id] = q

    # 파일 bytes 를 async 컨텍스트에서 읽고 나머지는 스레드에 위임
    raw_files: list[tuple[str, bytes]] = []
    for upload in files:
        raw = await upload.read()
        raw_files.append((upload.filename or "attachment", raw))

    loop = asyncio.get_running_loop()
    loop.run_in_executor(
        None,
        _run_simulation,
        session_id,
        topic,
        participants,
        rounds,
        raw_files,
    )

    return {"session_id": session_id}


@router.get("/stream/{session_id}")
async def stream_events(session_id: str) -> StreamingResponse:
    """SSE 스트림 — 전처리 및 시뮬레이션 이벤트를 실시간 전송합니다.

    SimpleQueue 를 50 ms 폴링으로 소비합니다.
    asyncio.Queue + call_soon_threadsafe 대신 이 방식이 Windows uvicorn 에서 안정적입니다.
    """
    q = _sessions.get(session_id)
    if q is None:
        raise HTTPException(status_code=404, detail="session not found")

    async def generate():
        try:
            while True:
                try:
                    event = q.get_nowait()
                except stdlib_queue.Empty:
                    await asyncio.sleep(0.05)  # 50 ms 폴링
                    continue

                if event is None:
                    yield 'data: {"type": "end"}\n\n'
                    break
                data = json.dumps(event, ensure_ascii=False)
                yield f"data: {data}\n\n"
        finally:
            _sessions.pop(session_id, None)

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )
